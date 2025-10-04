import os
import re
import pytesseract
from pptx import Presentation
from PIL import Image
import comtypes.client
from ai import file_description_and_keyfindings

# ------------------- Configure Tesseract -------------------
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

# ------------------- PII Masking -------------------
pii_patterns = {
    "EMPLOYEE_ID": r'\bEMP\d+\b',                        
    "EMAIL": r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b',
    "TOKEN_SERIAL": r'\bHT-\d+-[A-Z]+\b',                
    "FULL_NAME": r'\b(?:Mr|Mrs|Ms|Dr)\.?\s+[A-Z][a-z]+(?:\s[A-Z][a-z]+)+\b|\b[A-Z][a-z]+(?:\s[A-Z][a-z]+)+\b',
    "CREDIT_CARD": r"\b(?:\d[ -]*?){13,16}\b",
}


def mask_pii(text: str) -> str:
    """Mask PII while preserving table headers and context text."""
    if not text:
        return text

    masked_lines = []
    lines = text.splitlines()

    for line in lines:
        if re.search(r"\b(Revision|Full\s*Name|Email|Date|Changes)\b", line, re.IGNORECASE):
            masked_lines.append(line)
            continue


        masked_line = line
        masked_line = re.sub(r"[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}", "<EMAIL>", masked_line)
        masked_line = re.sub(r"\b(?:\+?\d{1,3}[-.\s]?)?(?:\(?\d{3}\)?[-.\s]?)?\d{3}[-.\s]?\d{4}\b", "<PHONE>", masked_line)

        # Full names: Two consecutive capitalized words (e.g., Sarah Thompson)
        masked_line = re.sub(
            r"\b([A-Z][a-z]+(?:\s+[A-Z][a-z]+)+)\b",
            lambda m: "<FULL_NAME>" if not re.search(r"(Changes|Network|Security|Policy|Diagram|Review)", m.group(1), re.IGNORECASE) else m.group(1),
            masked_line
        )

        masked_lines.append(masked_line)

    return "\n".join(masked_lines)


# ------------------- PPT Extraction -------------------
def extract_text_from_ppt(ppt_file: str) -> str:
    """Extract text from PPT/PPTX using python-pptx and OCR fallback."""
    text_content = ""

    # -------- Direct text extraction (text boxes + tables) --------
    try:
        prs = Presentation(ppt_file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):  # normal textbox
                    text_content += mask_pii(shape.text) + "\n"
                elif shape.has_table:  # table
                    for row in shape.table.rows:
                        row_text = "\t".join(mask_pii(cell.text.strip()) for cell in row.cells)
                        text_content += row_text + "\n"

    except Exception as e:
        print(f"Error extracting text directly: {e}")

    # -------- OCR fallback (for diagrams/images) --------
    try:
        out_dir = os.path.abspath("slides")
        os.makedirs(out_dir, exist_ok=True)

        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.Visible = 0  # set 1 if you want PowerPoint window to appear
        presentation = powerpoint.Presentations.Open(os.path.abspath(ppt_file))

        # Export slides as PNG images
        presentation.SaveAs(out_dir, 17)  # 17 = PNG
        presentation.Close()
        powerpoint.Quit()

        # OCR on exported slide images
        for img_file in sorted(os.listdir(out_dir)):
            if img_file.lower().endswith(".png"):
                img_path = os.path.join(out_dir, img_file)
                ocr_text = pytesseract.image_to_string(Image.open(img_path))
                if ocr_text.strip():
                    text_content += mask_pii(ocr_text) + "\n"

    except Exception as e:
        print(f"Error extracting OCR text: {e}")

    return text_content.strip()

def ppt_main(file_path: str):
    text_content = extract_text_from_ppt(file_path)
    if text_content:
        prompt = f"""
        Analyze the following cleansed presentation data (PPT/PPTX file). Generate a descriptive 
        title and a file caption of about 30 words that summarizes the data's content, 
        focusing on the security/IT-related context (e.g., Firewall rules, User Access Log, 
        Token Issuance, etc.).
        Cleansed Presentation Content:
        ---
        {text_content}
        """
        file_description,key_findings = file_description_and_keyfindings(prompt)
        return file_description,key_findings
    else:
        print("Failed to extract text from the PDF file.")




